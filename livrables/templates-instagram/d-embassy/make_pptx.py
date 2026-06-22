from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

BASE = os.path.dirname(os.path.abspath(__file__))

TEMPLATES = [
    ("template-1-conseil-du-jour.png",  "Conseil du jour"),
    ("template-2-question-du-jour.png", "Question du jour"),
    ("template-3-concept-du-jour.png",  "Concept du jour"),
]

SIZE = Inches(10)

prs = Presentation()
prs.slide_width  = SIZE
prs.slide_height = SIZE

blank_layout = prs.slide_layouts[6]

for filename, title in TEMPLATES:
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(BASE, filename)
    slide.shapes.add_picture(img_path, 0, 0, SIZE, SIZE)

out = os.path.join(BASE, "D Embassy Instagram Templates.pptx")
prs.save(out)
print(f"PPTX cree : {out}")
